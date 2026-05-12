"""PPT presentation export."""

from pathlib import Path
from datetime import datetime
from typing import Dict


def export_ppt(diff_result: Dict = None, enrich_result: Dict = None,
               storyline_result: Dict = None, output_dir: Path = None) -> Dict:
    """Generate a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                  "RNA-seq Analysis Report", 36, bold=True, color='1E3A8A')
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                  f"Generated: {datetime.now().strftime('%Y-%m-%d')}", 18, color='6B7280')

    # Slide 2: Data Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Data Overview", 28, bold=True, color='1E3A8A')
    content = "Expression data was analyzed for differential expression between treatment and control groups."
    if diff_result:
        content = f"Total genes: {diff_result.get('n_total', 'N/A')} | Up-regulated: {diff_result.get('n_up', 0)} | Down-regulated: {diff_result.get('n_down', 0)}\nMethod: {diff_result.get('method', 'N/A')}, |log2FC| > {diff_result.get('logfc_threshold', 1.0)}, FDR < {diff_result.get('fdr_threshold', 0.05)}"
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5), content, 16)

    # Slide 3: QC
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Quality Control", 28, bold=True, color='1E3A8A')
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5),
                  "QC checks performed: sample matching, missing values, duplicate genes, outliers, log2 transform need, group balance.", 16)

    # Slide 4: Differential Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Differential Expression Analysis", 28, bold=True, color='1E3A8A')
    top_genes_text = ""
    if diff_result:
        top = diff_result.get('top_genes', [])[:15]
        top_genes_text = "Top differentially expressed genes:\n\n" + "\n".join(
            f"  {g['gene']:15s}  log2FC={g.get('logfc', 0):+.2f}  FDR={g.get('fdr', 1):.1e}  [{g.get('direction', '')}]"
            for g in top
        )
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5), top_genes_text, 14)

    # Slide 5: Enrichment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Enrichment Analysis", 28, bold=True, color='1E3A8A')
    enrich_text = "No enrichment data available."
    if enrich_result:
        top = enrich_result.get('terms', [])[:10]
        enrich_text = "Top enriched terms:\n\n" + "\n".join(
            f"  {t['term'][:60]:60s}  FDR={t.get('fdr', 1):.1e}"
            for t in top
        )
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5), enrich_text, 12)

    # Slide 6: Key Genes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Key Gene Expression", 28, bold=True, color='1E3A8A')
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5),
                  "Key candidate genes identified from differential expression and enrichment analysis are highlighted for further validation.", 16)

    # Slide 7: Mechanism Hypothesis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Mechanism Hypothesis", 28, bold=True, color='1E3A8A')
    mech_text = "Run differential analysis and enrichment to generate mechanism hypotheses."
    if storyline_result and storyline_result.get('storylines'):
        s = storyline_result['storylines'][0]
        mech_text = s.get('mechanism', s.get('hypothesis', ''))
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5), mech_text[:500], 14)

    # Slide 8: Suggested Figures
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Suggested Figure Layout", 28, bold=True, color='1E3A8A')
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5),
                  "Figure 1: Volcano plot and heatmap\nFigure 2: GO/KEGG enrichment\nFigure 3: GSEA curves\nFigure 4: Key gene validation\nFigure 5: Mechanism model", 16)

    # Slide 9: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                  "Summary", 28, bold=True, color='1E3A8A')
    _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5),
                  "Analysis complete. Key findings, figures, and hypotheses generated. Ready for publication preparation.", 16)

    out_dir = Path(output_dir) if output_dir else Path('output')
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(str(path))

    return {"files": [{"name": "PPT Presentation", "path": str(path), "size": f"{path.stat().st_size/1024:.0f} KB"}]}


def _add_text_box(slide, left, top, width, height, text, size, bold=False, color='111827'):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*_hex_to_rgb(color))


def _hex_to_rgb(hex_color: str):
    h = hex_color.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
